import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0] # Title layout
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]
    title_box.text = title
    subtitle_box.text = subtitle

def add_content_slide(prs, title, content_lines):
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    title_box.text = title
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.text = content_lines[0]
    
    for line in content_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0 if not line.startswith("  ") else 1
        p.font.size = Pt(18)

def add_image_slide(prs, title, image_path, description_lines=None):
    slide_layout = prs.slide_layouts[5] # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    title_box.text = title
    
    if os.path.exists(image_path):
        # Add image
        # Center horizontally and place below title
        left = Inches(1.0)
        top = Inches(1.5)
        width = Inches(8.0)
        slide.shapes.add_picture(image_path, left, top, width=width)
    else:
        # Add a text box indicating missing image
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"Image {image_path} not found"

    if description_lines:
        left = Inches(0.5)
        top = Inches(5.7)
        width = Inches(9.0)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for line in description_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.level = 0 if not line.startswith("  ") else 1

prs = Presentation()

# Title Slide
add_title_slide(prs, "AI Cold Email Generator Pro", "Enterprise RAG + Local LLM Email & Resume Automation\n\nProject Overview & Architecture Explanation")

# Slide 1: Motive of the Project
add_content_slide(prs, "Motive of the Project", [
    "Why build the AI Cold Email Generator?",
    "  - Personalization at Scale: Writing tailored cold emails and resumes manually is extremely time-consuming.",
    "  - Data Privacy: Using public cloud LLMs for proprietary company data or personal PII poses severe security risks.",
    "  - Cost Efficiency: API costs for massive email campaigns can skyrocket. Local LLMs (Ollama) eliminate these ongoing expenses.",
    "  - Contextual Accuracy (RAG): Generic LLMs hallucinate. By indexing your specific PDFs and DOCXs (RAG), the AI writes emails containing factual, tailored information about your exact product or job description."
])

# Slide 2: Project Overview & Key Features
add_content_slide(prs, "Project Overview & Key Features", [
    "What is the system capable of?",
    "  - Project Management: Organize campaigns into isolated workspaces.",
    "  - Dynamic Knowledge Base: Upload PDF, DOCX, and TXT files per project.",
    "  - Retrieval-Augmented Generation (RAG): FAISS vector store combined with HuggingFace embeddings for precise context retrieval.",
    "  - Local AI Generation: Ollama (Gemma 2b model) processes context entirely on your machine to guarantee zero data leakage.",
    "  - Beautiful UI: Built with React 19, Tailwind CSS, Framer Motion, and Shadcn UI for a premium enterprise feel.",
    "  - Secure Authentication: JWT-based secure login and session management."
])

# Slide 3: High Level Architecture Diagram
add_image_slide(prs, "System Architecture", "architecture.png", [
    "This diagram illustrates the complete end-to-end flow of the application:",
    "  - Frontend (React) communicates with Backend (FastAPI) via JWT authenticated REST APIs.",
    "  - Backend acts as the orchestrator, communicating with the SQLite/PostgreSQL Database.",
    "  - The AI Engine is powered by LangChain, utilizing FAISS for RAG and Ollama (Gemma) for text generation."
])

# Slide 4: RAG Knowledge Ingestion Flow
add_image_slide(prs, "RAG Knowledge Ingestion Flow", "rag_flow.png", [
    "How does the AI securely 'read' and memorize documents?",
    "  - Documents (PDF, DOCX) are parsed via LangChain loaders.",
    "  - Text is chunked (1000 characters) to ensure the embedding model can process it.",
    "  - HuggingFace BGE Embeddings convert text chunks into mathematical vectors.",
    "  - Vectors are saved locally to a FAISS index tagged by project for future retrieval."
])

# Slide 5: AI Email Generation Flow
add_image_slide(prs, "AI Content Generation Flow", "generation_flow.png", [
    "How is a highly-tailored email generated?",
    "  - User requests an email by providing target company details and tone.",
    "  - The system retrieves the top 3 most relevant document chunks from the FAISS vector database.",
    "  - This contextual data is injected into a strict Prompt Template alongside the user's instructions.",
    "  - The local Ollama daemon (Gemma 2B) generates the email and streams the result back to the UI."
])

# Slide 6: Backend Execution Flow (Deep Dive)
add_content_slide(prs, "Backend Execution Flow (Deep Dive)", [
    "How does the backend process the generation request?",
    "  1. main.py intercepts the request (e.g. POST /api/emails/generate).",
    "  2. JWT Middleware validates the user's token.",
    "  3. Pydantic strictly validates the incoming JSON payload.",
    "  4. The AIService queries the local FAISS index for RAG context.",
    "  5. LangChain orchestrates the prompt assembly and local LLM execution.",
    "  6. The generated content is saved to the DB and returned to the client."
])

# Slide 7: Frontend Experience Flow (Deep Dive)
add_content_slide(prs, "Frontend Experience Flow (Deep Dive)", [
    "How does the UI deliver a seamless user experience?",
    "  1. Zustand manages global application state (like User Auth).",
    "  2. React Hook Form + Zod provide instant, rigorous form validation.",
    "  3. TanStack React Query handles data fetching, caching, and background syncing.",
    "  4. Axios interceptors automatically attach JWT tokens to all requests.",
    "  5. Framer Motion provides smooth skeleton loaders during the AI generation process."
])

# Slide 8: Conclusion & Q&A
add_content_slide(prs, "Conclusion", [
    "Thank You for your time!",
    "",
    "Summary:",
    "  - A complete enterprise solution for automated cold outreach.",
    "  - 100% Data Privacy via Local AI.",
    "  - Highly tailored and factual emails via RAG.",
    "",
    "Any Questions?"
])

prs.save("Project_Explanation_Presentation.pptx")
print("Presentation generated successfully!")
